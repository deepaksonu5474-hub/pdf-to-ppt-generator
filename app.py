import streamlit as st
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import pytesseract
import io
import os

st.set_page_config(
    page_title="PDF to PPT Generator",
    page_icon="📄",
    layout="wide"
)

st.title("📄 PDF to Editable PPT Generator")
st.write("Convert PDF and Images into editable PowerPoint.")

st.divider()

uploaded_file = st.file_uploader(
    "Upload your file",
    type=[
        "pdf",
        "png",
        "jpg",
        "jpeg",
        "webp"
    ]
)

if uploaded_file:

    file_name = uploaded_file.name
    file_ext = os.path.splitext(file_name)[1].lower()

    st.success(f"File selected: {file_name}")

    if st.button(
        "🚀 Generate Editable PowerPoint",
        type="primary"
    ):

        with st.spinner(
            "Reading file and creating editable PPT..."
        ):

            file_bytes = uploaded_file.read()

            prs = Presentation()

            # 16:9 widescreen
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # ============================================
            # IMAGE → OCR → EDITABLE TEXT
            # ============================================

            if file_ext in [
                ".png",
                ".jpg",
                ".jpeg",
                ".webp"
            ]:

                image = Image.open(
                    io.BytesIO(file_bytes)
                )

                # Convert to RGB
                image = image.convert("RGB")

                slide = prs.slides.add_slide(
                    prs.slide_layouts[6]
                )

                # OCR
                ocr_data = pytesseract.image_to_data(
                    image,
                    lang="hin+eng",
                    output_type=pytesseract.Output.DICT
                )

                img_width, img_height = image.size

                for i in range(
                    len(ocr_data["text"])
                ):

                    text = ocr_data["text"][i].strip()

                    if not text:
                        continue

                    confidence = float(
                        ocr_data["conf"][i]
                    )

                    if confidence < 30:
                        continue

                    x = ocr_data["left"][i]
                    y = ocr_data["top"][i]
                    w = ocr_data["width"][i]
                    h = ocr_data["height"][i]

                    left = Inches(
                        (x / img_width) * 13.333
                    )

                    top = Inches(
                        (y / img_height) * 7.5
                    )

                    width = Inches(
                        (w / img_width) * 13.333
                    )

                    height = Inches(
                        (h / img_height) * 7.5
                    )

                    textbox = slide.shapes.add_textbox(
                        left,
                        top,
                        width,
                        height
                    )

                    text_frame = textbox.text_frame
                    text_frame.clear()

                    paragraph = (
                        text_frame.paragraphs[0]
                    )

                    run = paragraph.add_run()

                    run.text = text

                    run.font.size = Pt(16)

                    paragraph.alignment = (
                        PP_ALIGN.LEFT
                    )

            # ============================================
            # PDF → EDITABLE TEXT
            # ============================================

            elif file_ext == ".pdf":

                pdf = fitz.open(
                    stream=file_bytes,
                    filetype="pdf"
                )

                for page in pdf:

                    slide = prs.slides.add_slide(
                        prs.slide_layouts[6]
                    )

                    page_width = page.rect.width
                    page_height = page.rect.height

                    blocks = page.get_text(
                        "blocks"
                    )

                    for block in blocks:

                        x0, y0, x1, y1, text = block[:5]

                        text = text.strip()

                        if not text:
                            continue

                        left = Inches(
                            (x0 / page_width) * 13.333
                        )

                        top = Inches(
                            (y0 / page_height) * 7.5
                        )

                        width = Inches(
                            ((x1 - x0) / page_width) * 13.333
                        )

                        height = Inches(
                            ((y1 - y0) / page_height) * 7.5
                        )

                        textbox = slide.shapes.add_textbox(
                            left,
                            top,
                            width,
                            height
                        )

                        text_frame = textbox.text_frame
                        text_frame.clear()

                        paragraph = (
                            text_frame.paragraphs[0]
                        )

                        run = paragraph.add_run()

                        run.text = text

                        run.font.size = Pt(16)

                        paragraph.alignment = (
                            PP_ALIGN.LEFT
                        )

            # ============================================
            # SAVE PPT
            # ============================================

            output = io.BytesIO()

            prs.save(output)

            output.seek(0)

        st.success(
            "✅ Editable PowerPoint created!"
        )

        st.download_button(
            label="⬇️ Download Editable PowerPoint",
            data=output,
            file_name="editable_presentation.pptx",
            mime=(
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            )
        )
